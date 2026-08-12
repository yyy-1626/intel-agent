"""Office 文档适配器 — 提取 Word(.docx)、Excel(.xlsx)、PowerPoint(.pptx) 文本"""

from __future__ import annotations

import io
import logging
from typing import Optional

import requests

from .base import ContentAdapter

logger = logging.getLogger(__name__)

# Office MIME 类型
OFFICE_MIME = {
    "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
    "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
    "application/vnd.openxmlformats-officedocument.presentationml.presentation",
    "application/msword",
    "application/vnd.ms-excel",
    "application/vnd.ms-powerpoint",
}

# URL 后缀映射
OFFICE_SUFFIXES = {".docx", ".xlsx", ".pptx", ".doc", ".xls", ".ppt"}


class OfficeAdapter(ContentAdapter):
    """从 Office 文档中提取文本内容。"""

    def can_handle(self, content_type: str, url: str) -> bool:
        ct = content_type.lower().split(";")[0].strip()
        if ct in OFFICE_MIME:
            return True
        url_lower = url.lower().split("?")[0]
        return any(url_lower.endswith(s) for s in OFFICE_SUFFIXES)

    def extract(self, url: str, response: Optional[requests.Response] = None) -> str:
        if response is None:
            response = requests.get(url, timeout=30)
            response.raise_for_status()

        content = response.content
        url_lower = url.lower()

        if url_lower.endswith((".docx", ".doc")) or self._is_docx(response):
            return self._extract_docx(content)
        elif url_lower.endswith((".xlsx", ".xls")):
            return self._extract_xlsx(content)
        elif url_lower.endswith((".pptx", ".ppt")):
            return self._extract_pptx(content)
        else:
            # 根据 Content-Type 兜底判断
            ct = response.headers.get("Content-Type", "").lower()
            if "word" in ct:
                return self._extract_docx(content)
            elif "spreadsheet" in ct or "excel" in ct:
                return self._extract_xlsx(content)
            elif "presentation" in ct or "powerpoint" in ct:
                return self._extract_pptx(content)
            return ""

    def _is_docx(self, response) -> bool:
        ct = response.headers.get("Content-Type", "").lower()
        return "word" in ct or "msword" in ct

    def _extract_docx(self, content: bytes) -> str:
        try:
            from docx import Document
            doc = Document(io.BytesIO(content))
            paragraphs = [p.text for p in doc.paragraphs if p.text.strip()]
            text = "\n".join(paragraphs)
            logger.info("OfficeAdapter(docx): %d 段，%d 字", len(paragraphs), len(text))
            return text
        except ImportError:
            logger.warning("OfficeAdapter: python-docx 未安装")
            return ""
        except Exception as e:
            logger.warning("OfficeAdapter(docx) 失败: %s", e)
            return ""

    def _extract_xlsx(self, content: bytes) -> str:
        try:
            from openpyxl import load_workbook
            wb = load_workbook(io.BytesIO(content), read_only=True, data_only=True)
            all_rows: list[str] = []
            for sheet_name in wb.sheetnames:
                ws = wb[sheet_name]
                sheet_data = [f"[工作表: {sheet_name}]"]
                for row in ws.iter_rows(values_only=True):
                    row_text = "\t".join(str(c) if c is not None else "" for c in row)
                    if row_text.strip():
                        sheet_data.append(row_text)
                all_rows.append("\n".join(sheet_data))
            wb.close()
            text = "\n\n".join(all_rows)
            logger.info("OfficeAdapter(xlsx): %d 表，%d 字", len(wb.sheetnames), len(text))
            return text
        except ImportError:
            logger.warning("OfficeAdapter: openpyxl 未安装")
            return ""
        except Exception as e:
            logger.warning("OfficeAdapter(xlsx) 失败: %s", e)
            return ""

    def _extract_pptx(self, content: bytes) -> str:
        try:
            from pptx import Presentation
            prs = Presentation(io.BytesIO(content))
            slides_text: list[str] = []
            for i, slide in enumerate(prs.slides):
                slide_lines = [f"[幻灯片 {i + 1}]"]
                for shape in slide.shapes:
                    if shape.has_text_frame:
                        for para in shape.text_frame.paragraphs:
                            t = para.text.strip()
                            if t:
                                slide_lines.append(t)
                if len(slide_lines) > 1:
                    slides_text.append("\n".join(slide_lines))
            text = "\n\n".join(slides_text)
            logger.info("OfficeAdapter(pptx): %d 页有文本，%d 字", len(slides_text), len(text))
            return text
        except ImportError:
            logger.warning("OfficeAdapter: python-pptx 未安装")
            return ""
        except Exception as e:
            logger.warning("OfficeAdapter(pptx) 失败: %s", e)
            return ""
